"""업로드 파일 → 평문 텍스트 추출.

지원: txt / html / pptx / docx / pdf.
각 형식의 결과는 "위치 표시(location)"와 함께 (text, location) 조각 리스트로 돌려줘서,
근거 표시에 슬라이드/페이지/섹션 번호 등을 쓸 수 있게 한다.
"""
import io

SUPPORTED = {"txt", "md", "html", "htm", "pptx", "docx", "pdf"}


def _ext(filename: str) -> str:
    return filename.lower().rsplit(".", 1)[-1] if "." in filename else ""


def parse_file(filename: str, data: bytes) -> list[tuple[str, str]]:
    """파일을 (text, location) 조각 리스트로 파싱."""
    ext = _ext(filename)
    if ext in ("txt", "md"):
        return [(data.decode("utf-8", errors="ignore"), "본문")]
    if ext in ("html", "htm"):
        return [(_parse_html(data), "본문")]
    if ext == "pptx":
        return _parse_pptx(data)
    if ext == "docx":
        return [(_parse_docx(data), "본문")]
    if ext == "pdf":
        return _parse_pdf(data)
    raise ValueError(f"아직 지원하지 않는 형식입니다: .{ext} (지원: {', '.join(sorted(SUPPORTED))})")


def _parse_html(data: bytes) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(data, "lxml")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text("\n")


def _parse_pptx(data: bytes) -> list[tuple[str, str]]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    out: list[tuple[str, str]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        joined = "\n".join(t for t in texts if t.strip())
        if joined.strip():
            out.append((joined, f"슬라이드 {idx}"))
    return out


def _parse_docx(data: bytes) -> str:
    """DOCX → 평문. 본문 읽기 순서(문단·표 교차)를 유지한다.

    NOTE: `doc.paragraphs` 전부 + `doc.tables` 전부 순으로 붙이면
    '실적 N' 제목과 바로 아래 발주처 표가 떨어져, RAG가 3번째 실적 메타를
    제목과 연결하지 못하는 문제가 난다.
    """
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            text = Paragraph(child, doc).text
            if text.strip():
                parts.append(text)
        elif child.tag == qn("w:tbl"):
            table = Table(child, doc)
            for row in table.rows:
                # 병합 셀 중복 텍스트는 한 번만
                cells: list[str] = []
                seen_tc: set[int] = set()
                for cell in row.cells:
                    tc_id = id(cell._tc)
                    if tc_id in seen_tc:
                        continue
                    seen_tc.add(tc_id)
                    t = cell.text.strip()
                    if t:
                        cells.append(t)
                if cells:
                    parts.append(" | ".join(cells))
    return "\n".join(parts)


def _parse_pdf(data: bytes) -> list[tuple[str, str]]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    out: list[tuple[str, str]] = []
    for idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            out.append((text, f"페이지 {idx}"))
    return out
